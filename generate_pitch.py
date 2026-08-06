import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

def add_screenshot_placeholder(slide):
    # Bottom half placeholder: left=1, top=4.5, width=8, height=2.5
    img_path = r"C:\Users\Driftingdane\.gemini\antigravity-ide\brain\ef54668a-8a99-4bfb-ab9e-129c8613dd98\dummy_placeholder_1783364209416.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(4.5), width=Inches(8), height=Inches(2.5))

def create_presentation():
    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    base_dir = os.path.dirname(os.path.abspath(__file__))

    # Slide 1: Title
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Fra OneNote til Q&A-viden"
    slide.placeholders[1].text = "Morten Noerregaard (OOZ)"

    # Helper to setup text body for top-stacked layout
    def setup_top_text(slide):
        body_shape = slide.placeholders[1]
        body_shape.left = Inches(0.5)
        body_shape.top = Inches(1.5)
        body_shape.width = Inches(9)
        body_shape.height = Inches(3.0)
        tf = body_shape.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        return tf

    # Slide 2: Hvad savner vi fra OneNote?
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Hvad savner vi fra OneNote?"
    tf = setup_top_text(slide)
    tf.text = "En fremtidsikret videnbase som kan viderudvikles."
    for point in ["Struktureret data", "Kategorisering", "Søgning & Deling"]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
    add_screenshot_placeholder(slide)

    # Slide 3: Hvad er anderledes? Søgning der virker
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Hvad er anderledes? Søgning der virker"
    tf = setup_top_text(slide)
    tf.text = "Skriv i søgefeltet, og listen filtrerer med det samme."
    for point in ["Du kan sortere efter dato, bruger og filtrere på kategorier i sidepanelet.", "Hoppe direkte mellem dine søgeord i teksten – via quick links."]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
    add_screenshot_placeholder(slide)

    # Slide 4: Automatiske tags & Linking
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Automatiske tags & Linking"
    tf = setup_top_text(slide)
    tf.text = "Automatiske tags:"
    p = tf.add_paragraph(); p.text = "Systemet tilknytter tags, så du kan klikke dig frem til sager med samme problemtype."; p.level = 1
    p = tf.add_paragraph(); p.text = "Linking:"; p.level = 0
    p = tf.add_paragraph(); p.text = "Hjælp dine kollegaer med viden med direkte links uden at skulle søge."; p.level = 1
    add_screenshot_placeholder(slide)

    # Slide 5: Skærmbilleder & Sikkerhed
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Skærmbilleder & Sikkerhed"
    tf = setup_top_text(slide)
    tf.text = "Automatisk sløring af skærmbilleder:"
    p = tf.add_paragraph(); p.text = "Fjerner CPR-nummer og bank oplysninger automatisk inden billedet gemmes."; p.level = 1
    p = tf.add_paragraph(); p.text = "Microsoft-login:"; p.level = 0
    p = tf.add_paragraph(); p.text = "Koblet direkte på virksomhedens Microsoft-miljø (automatisk login)."; p.level = 1
    add_screenshot_placeholder(slide)

    # Slide 6: Konklusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "En skalerbar løsning"
    tf = setup_top_text(slide)
    tf.text = "Sikkerhed og overholdelse af GDPR:"
    p = tf.add_paragraph(); p.text = "Fjerner risikoen for at sløret data glemmes i uendelige notesbøger."; p.level = 1
    p = tf.add_paragraph(); p.text = "Struktureret format:"; p.level = 0
    p = tf.add_paragraph(); p.text = "Uorganiseret tekst bliver fortid – gør det muligt at filtrere effektivt."; p.level = 1
    p = tf.add_paragraph(); p.text = "Med Q&A-vidensbase bygger vi et skalerbart system."; p.level = 0
    
    output_path = os.path.join(base_dir, "Fra_OneNote_Til_QA_Template.pptx")
    prs.save(output_path)
    print(f"Presentation template generated successfully at: {output_path}")

if __name__ == '__main__':
    create_presentation()
